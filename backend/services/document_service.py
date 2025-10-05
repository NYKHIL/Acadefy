"""
Acadefy Document Service
Handles document upload, processing, and knowledge extraction
"""

import os
import requests
import logging
from typing import Dict, List, Any, Optional
import re
from urllib.parse import urlparse
import json
from werkzeug.utils import secure_filename
from werkzeug.utils import secure_filename

# Import document processing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentService:
    """
    Service for processing documents and extracting knowledge
    """
    
    def __init__(self):
        self.documents_dir = "documents"
        self.uploads_dir = os.path.join(self.documents_dir, "uploads")
        self.knowledge_base = {}
        self.allowed_extensions = {'pdf', 'docx', 'pptx', 'txt'}
        self.ensure_documents_directory()
        self.load_knowledge_base()
    
    def ensure_documents_directory(self):
        """Create documents directory if it doesn't exist"""
        if not os.path.exists(self.documents_dir):
            os.makedirs(self.documents_dir)
        if not os.path.exists(self.uploads_dir):
            os.makedirs(self.uploads_dir)
    
    def add_document_from_url(self, url: str, title: str = None) -> Dict[str, Any]:
        """
        Add a document from URL (supports text content, PDFs, etc.)
        """
        try:
            # Validate URL
            parsed_url = urlparse(url)
            if not parsed_url.scheme or not parsed_url.netloc:
                return {"success": False, "error": "Invalid URL format"}
            
            # Download content
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Determine content type
            content_type = response.headers.get('content-type', '').lower()
            
            if 'text' in content_type or 'html' in content_type:
                content = self._extract_text_from_html(response.text)
            elif 'pdf' in content_type:
                content = self._extract_text_from_pdf(response.content)
            else:
                content = response.text
            
            # Generate document ID
            doc_id = self._generate_doc_id(url, title)
            
            # Store document
            document = {
                "id": doc_id,
                "title": title or self._extract_title_from_url(url),
                "url": url,
                "content": content,
                "content_type": content_type,
                "chunks": self._chunk_content(content),
                "keywords": self._extract_keywords(content)
            }
            
            self.knowledge_base[doc_id] = document
            self.save_knowledge_base()
            
            logger.info(f"Document added successfully: {doc_id}")
            return {
                "success": True, 
                "document_id": doc_id,
                "title": document["title"],
                "chunks_count": len(document["chunks"])
            }
            
        except requests.RequestException as e:
            logger.error(f"Error downloading document: {str(e)}")
            return {"success": False, "error": f"Failed to download: {str(e)}"}
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            return {"success": False, "error": f"Processing failed: {str(e)}"}
    
    def add_document_from_text(self, content: str, title: str, source: str = "manual") -> Dict[str, Any]:
        """
        Add a document from direct text input
        """
        try:
            doc_id = self._generate_doc_id(source, title)
            
            document = {
                "id": doc_id,
                "title": title,
                "url": source,
                "content": content,
                "content_type": "text/plain",
                "chunks": self._chunk_content(content),
                "keywords": self._extract_keywords(content)
            }
            
            self.knowledge_base[doc_id] = document
            self.save_knowledge_base()
            
            return {
                "success": True,
                "document_id": doc_id,
                "title": title,
                "chunks_count": len(document["chunks"])
            }
            
        except Exception as e:
            logger.error(f"Error adding text document: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def add_document_from_file(self, file, title: str = None) -> Dict[str, Any]:
        """
        Add a document from uploaded file (PDF, DOCX, PPTX, TXT)
        """
        try:
            if not file or not file.filename:
                return {"success": False, "error": "No file provided"}
            
            # Check file extension
            filename = secure_filename(file.filename)
            file_ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
            
            if file_ext not in self.allowed_extensions:
                return {"success": False, "error": f"Unsupported file type. Allowed: {', '.join(self.allowed_extensions)}"}
            
            # Save file
            file_path = os.path.join(self.uploads_dir, filename)
            file.save(file_path)
            
            # Extract content based on file type
            if file_ext == 'pdf':
                content = self._extract_text_from_pdf_file(file_path)
            elif file_ext == 'docx':
                content = self._extract_text_from_docx_file(file_path)
            elif file_ext == 'pptx':
                content = self._extract_text_from_pptx_file(file_path)
            elif file_ext == 'txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            else:
                return {"success": False, "error": "Unsupported file type"}
            
            if not content.strip():
                return {"success": False, "error": "No text content found in file"}
            
            # Generate document ID
            doc_id = self._generate_doc_id(filename, title or filename)
            
            # Store document
            document = {
                "id": doc_id,
                "title": title or filename,
                "url": f"file://{filename}",
                "file_path": file_path,
                "content": content,
                "content_type": f"application/{file_ext}",
                "chunks": self._chunk_content(content),
                "keywords": self._extract_keywords(content),
                "file_size": os.path.getsize(file_path)
            }
            
            self.knowledge_base[doc_id] = document
            self.save_knowledge_base()
            
            logger.info(f"File document added successfully: {doc_id}")
            return {
                "success": True,
                "document_id": doc_id,
                "title": document["title"],
                "chunks_count": len(document["chunks"]),
                "file_size": document["file_size"]
            }
            
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            return {"success": False, "error": f"File processing failed: {str(e)}"}
    
    def search_documents(self, query: str, max_results: int = 3) -> List[Dict[str, Any]]:
        """
        Search through documents for relevant content with improved matching
        """
        query_lower = query.lower()
        query_words = query_lower.split()
        results = []
        
        logger.info(f"Searching documents for query: '{query}' (words: {query_words})")
        logger.info(f"Total documents in knowledge base: {len(self.knowledge_base)}")
        
        for doc_id, document in self.knowledge_base.items():
            relevance_score = 0
            matching_chunks = []
            
            logger.info(f"Checking document: {document['title']}")
            
            # Check title relevance (exact phrase)
            if query_lower in document["title"].lower():
                relevance_score += 10
                logger.info(f"  Title match found: +10 points")
            
            # Check title relevance (individual words)
            for word in query_words:
                if word in document["title"].lower():
                    relevance_score += 3
                    logger.info(f"  Title word match '{word}': +3 points")
            
            # Check keyword relevance (both ways)
            for keyword in document["keywords"]:
                keyword_lower = keyword.lower()
                # Check if query contains keyword
                if keyword_lower in query_lower:
                    relevance_score += 5
                    logger.info(f"  Keyword match '{keyword}': +5 points")
                # Check if keyword contains any query word
                for word in query_words:
                    if word in keyword_lower and len(word) > 2:
                        relevance_score += 2
                        logger.info(f"  Keyword partial match '{keyword}' contains '{word}': +2 points")
            
            # Check content chunks (improved matching)
            for i, chunk in enumerate(document["chunks"]):
                chunk_lower = chunk.lower()
                chunk_score = 0
                
                # Exact phrase match
                if query_lower in chunk_lower:
                    chunk_score += 5
                    logger.info(f"  Chunk {i} exact phrase match: +5 points")
                
                # Individual word matches
                word_matches = 0
                for word in query_words:
                    if len(word) > 2 and word in chunk_lower:  # Skip very short words
                        word_matches += 1
                
                if word_matches > 0:
                    chunk_score += word_matches * 2
                    logger.info(f"  Chunk {i} word matches ({word_matches}): +{word_matches * 2} points")
                
                if chunk_score > 0:
                    relevance_score += chunk_score
                    matching_chunks.append({
                        "chunk_index": i,
                        "content": chunk,
                        "relevance": chunk_score,
                        "word_matches": word_matches
                    })
            
            logger.info(f"  Total relevance score for {document['title']}: {relevance_score}")
            
            # Lower threshold for inclusion
            if relevance_score > 0:
                results.append({
                    "document_id": doc_id,
                    "title": document["title"],
                    "url": document["url"],
                    "relevance_score": relevance_score,
                    "matching_chunks": sorted(matching_chunks, 
                                            key=lambda x: x["relevance"], 
                                            reverse=True)[:3]  # Top 3 chunks
                })
        
        # Sort by relevance and return top results
        results.sort(key=lambda x: x["relevance_score"], reverse=True)
        logger.info(f"Search completed. Found {len(results)} relevant documents")
        
        return results[:max_results]
    
    def get_document_context(self, query: str) -> str:
        """
        Get relevant document context for a query with improved formatting
        """
        search_results = self.search_documents(query)
        
        logger.info(f"Getting context for query '{query}': found {len(search_results)} relevant documents")
        
        if not search_results:
            # If no specific matches, try to return some content from any document
            if self.knowledge_base:
                logger.info("No specific matches found, returning sample content from available documents")
                doc = next(iter(self.knowledge_base.values()))
                return f"**From: {doc['title']}**\n{doc['chunks'][0] if doc['chunks'] else 'No content available'}"
            return ""
        
        context_parts = []
        for result in search_results:
            context_parts.append(f"**Document: {result['title']} (Relevance: {result['relevance_score']})**")
            
            if result["matching_chunks"]:
                for chunk in result["matching_chunks"]:
                    context_parts.append(f"Content: {chunk['content']}")
            else:
                # If no matching chunks, include first chunk as context
                doc = self.knowledge_base.get(result['document_id'])
                if doc and doc['chunks']:
                    context_parts.append(f"Content: {doc['chunks'][0]}")
            
            context_parts.append("")  # Empty line between documents
        
        context = "\n".join(context_parts)
        logger.info(f"Generated context length: {len(context)} characters")
        
        return context
    
    def list_documents(self) -> List[Dict[str, Any]]:
        """
        List all documents in the knowledge base
        """
        return [
            {
                "id": doc_id,
                "title": doc["title"],
                "url": doc["url"],
                "content_type": doc["content_type"],
                "chunks_count": len(doc["chunks"]),
                "keywords_count": len(doc["keywords"])
            }
            for doc_id, doc in self.knowledge_base.items()
        ]
    
    def remove_document(self, doc_id: str) -> bool:
        """
        Remove a document from the knowledge base
        """
        if doc_id in self.knowledge_base:
            del self.knowledge_base[doc_id]
            self.save_knowledge_base()
            return True
        return False
    
    # Helper methods
    def _generate_doc_id(self, source: str, title: str) -> str:
        """Generate unique document ID"""
        import hashlib
        import time
        content = f"{source}_{title}_{time.time()}"
        return hashlib.md5(content.encode()).hexdigest()[:12]
    
    def _extract_title_from_url(self, url: str) -> str:
        """Extract title from URL"""
        parsed = urlparse(url)
        filename = os.path.basename(parsed.path)
        return filename if filename else parsed.netloc
    
    def _extract_text_from_html(self, html_content: str) -> str:
        """Extract text from HTML content"""
        # Simple HTML tag removal (for basic cases)
        import re
        text = re.sub(r'<[^>]+>', '', html_content)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()
    
    def _extract_text_from_pdf(self, pdf_content: bytes) -> str:
        """Extract text from PDF content"""
        if not PDF_AVAILABLE:
            return "PDF processing not available. Install PyPDF2: pip install PyPDF2"
        
        try:
            import io
            pdf_file = io.BytesIO(pdf_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"
    
    def _extract_text_from_pdf_file(self, file_path: str) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            return "PDF processing not available. Install PyPDF2: pip install PyPDF2"
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return f"Error extracting PDF text: {str(e)}"
    
    def _extract_text_from_docx_file(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            return "DOCX processing not available. Install python-docx: pip install python-docx"
        
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            return f"Error extracting DOCX text: {str(e)}"
    
    def _extract_text_from_pptx_file(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            return "PPTX processing not available. Install python-pptx: pip install python-pptx"
        
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            return f"Error extracting PPTX text: {str(e)}"
    
    def _chunk_content(self, content: str, chunk_size: int = 300) -> List[str]:
        """Split content into manageable chunks by sentences and paragraphs"""
        # First, split by paragraphs
        paragraphs = content.split('\n\n')
        chunks = []
        
        current_chunk = ""
        current_word_count = 0
        
        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
                
            # Split paragraph into sentences
            sentences = re.split(r'[.!?]+', paragraph)
            
            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue
                    
                sentence_words = len(sentence.split())
                
                # If adding this sentence would exceed chunk size, save current chunk
                if current_word_count + sentence_words > chunk_size and current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = sentence + ". "
                    current_word_count = sentence_words
                else:
                    current_chunk += sentence + ". "
                    current_word_count += sentence_words
        
        # Add the last chunk if it has content
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # If no chunks were created, split by words as fallback
        if not chunks:
            words = content.split()
            for i in range(0, len(words), chunk_size):
                chunk = ' '.join(words[i:i + chunk_size])
                if chunk.strip():
                    chunks.append(chunk.strip())
        
        return chunks
    
    def _extract_keywords(self, content: str, max_keywords: int = 20) -> List[str]:
        """Extract keywords from content"""
        # Simple keyword extraction
        words = re.findall(r'\b[a-zA-Z]{4,}\b', content.lower())
        
        # Remove common stop words
        stop_words = {'this', 'that', 'with', 'have', 'will', 'from', 'they', 'been', 'said', 'each', 'which', 'their', 'time', 'about', 'would', 'there', 'could', 'other', 'more', 'very', 'what', 'know', 'just', 'first', 'into', 'over', 'think', 'also', 'your', 'work', 'life', 'only', 'can', 'still', 'should', 'after', 'being', 'now', 'made', 'before', 'here', 'through', 'when', 'where', 'much', 'some', 'these', 'many', 'then', 'them', 'well', 'were'}
        
        filtered_words = [word for word in words if word not in stop_words and len(word) > 3]
        
        # Count frequency and return most common
        from collections import Counter
        word_counts = Counter(filtered_words)
        return [word for word, count in word_counts.most_common(max_keywords)]
    
    def _calculate_chunk_relevance(self, query: str, chunk: str) -> float:
        """Calculate relevance score for a chunk"""
        query_words = query.split()
        chunk_words = chunk.split()
        
        matches = sum(1 for word in query_words if word in chunk_words)
        return matches / len(query_words) if query_words else 0
    
    def save_knowledge_base(self):
        """Save knowledge base to file"""
        try:
            kb_file = os.path.join(self.documents_dir, "knowledge_base.json")
            with open(kb_file, 'w', encoding='utf-8') as f:
                json.dump(self.knowledge_base, f, indent=2, ensure_ascii=False)
        except Exception as e:
            logger.error(f"Error saving knowledge base: {str(e)}")
    
    def load_knowledge_base(self):
        """Load knowledge base from file"""
        try:
            kb_file = os.path.join(self.documents_dir, "knowledge_base.json")
            if os.path.exists(kb_file):
                with open(kb_file, 'r', encoding='utf-8') as f:
                    self.knowledge_base = json.load(f)
                logger.info(f"Loaded {len(self.knowledge_base)} documents from knowledge base")
        except Exception as e:
            logger.error(f"Error loading knowledge base: {str(e)}")
            self.knowledge_base = {}
    
    def _extract_text_from_pdf_file(self, file_path: str) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            return "PDF processing requires PyPDF2 library. Install with: pip install PyPDF2"
        
        try:
            import PyPDF2
            text_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
            
            return '\n'.join(text_content)
        except Exception as e:
            logger.error(f"Error extracting PDF text from file: {str(e)}")
            return f"Error processing PDF file: {str(e)}"
    
    def _extract_text_from_docx_file(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            return "DOCX processing requires python-docx library. Install with: pip install python-docx"
        
        try:
            from docx import Document
            doc = Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            
            return '\n'.join(text_content)
        except Exception as e:
            logger.error(f"Error extracting DOCX text from file: {str(e)}")
            return f"Error processing DOCX file: {str(e)}"
    
    def _extract_text_from_pptx_file(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            return "PPTX processing requires python-pptx library. Install with: pip install python-pptx"
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            return '\n'.join(text_content)
        except Exception as e:
            logger.error(f"Error extracting PPTX text from file: {str(e)}")
            return f"Error processing PPTX file: {str(e)}"
    
